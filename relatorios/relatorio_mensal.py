#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
relatorio_mensal.py — gera o relatório mensal "Demonstrações Financeiras" em
PowerPoint EDITÁVEL (.pptx) a partir dos dados reais do cockpit (cockpit_ref).

Feito com python-pptx → gráficos NATIVOS do PowerPoint (o cliente edita cores,
dados, adiciona notas/slides). Roda no container do FastAPI (server-side) e é
disparado pelo cockpit (GET /api/relatorio/{slug}?ano=&mes=) ou pela CLI.

Reutiliza a MESMA lógica de dados dos endpoints (api_cockpit) — nada hardcoded.
Cobre grupo consolidado E cada empresa. Marca meses PROJETADOS (> realizado_ate).

Uso CLI:  python relatorios/relatorio_mensal.py grupo 2026 6
"""
import io, os, sys, datetime as dt

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "ia"))
import api_cockpit as A                       # reaproveita _conn, _dre_*, constantes, EMPRESAS

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ── Identidade REF (tema claro do cockpit) ───────────────────────────────────
INK   = RGBColor(0x1C, 0x1C, 0x1C)
PAPER = RGBColor(0xF9, 0xF8, 0xF6)
GRAY  = RGBColor(0x81, 0x80, 0x7C)
LINE  = RGBColor(0xE6, 0xE3, 0xDC)
ACCENT= RGBColor(0xD9, 0xDA, 0x00)           # amarelo REF
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED   = RGBColor(0xE5, 0x48, 0x4D)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Calibri"                            # seguro (render fiel + ships c/ Office)

MESES = ["Jan", "Fev", "Mar", "Abr", "Mai", "Jun",
         "Jul", "Ago", "Set", "Out", "Nov", "Dez"]
MES_EXT = ["Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
           "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"]

W, H = Inches(13.333), Inches(7.5)           # 16:9 wide


# ── helpers de formatação (pt-BR) ────────────────────────────────────────────
def _mi(v):
    if v is None:
        return "—"
    return ("R$ %.1f mi" % (v / 1e6)).replace(".", ",")


def _pct(v, nd=1):
    if v is None:
        return "—"
    return (f"%.{nd}f%%" % v).replace(".", ",")


# ── helpers de layout ────────────────────────────────────────────────────────
def _fundo(slide, cor=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = cor


def _cx(slide, x, y, w, h, texto, size, cor=INK, bold=False, align=PP_ALIGN.LEFT,
        font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    linhas = texto.split("\n")
    for i, ln in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = cor
    return tb


def _retangulo(slide, x, y, w, h, fill=WHITE, linha=LINE):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = linha; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def _estilo_grafico(chart, cores, com_legenda=False, rotulos=False,
                    rotulo_pos=XL_LABEL_POSITION.OUTSIDE_END, num_fmt=None):
    chart.has_title = False
    chart.font.name = FONT; chart.font.size = Pt(10); chart.font.color.rgb = GRAY
    if com_legenda:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
    else:
        chart.has_legend = False
    for i, plot in enumerate(chart.plots):
        plot.gap_width = 60
        for j, serie in enumerate(plot.series):
            serie.format.fill.solid()
            serie.format.fill.fore_color.rgb = cores[j % len(cores)]
        if rotulos:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(9); dl.font.name = FONT; dl.font.color.rgb = INK
            dl.position = rotulo_pos
            if num_fmt:
                dl.number_format = num_fmt; dl.number_format_is_linked = False


# ── slides ───────────────────────────────────────────────────────────────────
def _slide_capa(prs, titulo_emp, ano, mes, realizado_ate):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _fundo(s, INK)
    # bloco amarelo de marca (motivo visual: quadrado accent, sem faixas)
    from pptx.enum.shapes import MSO_SHAPE
    q = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.6), Inches(0.55), Inches(0.55))
    q.fill.solid(); q.fill.fore_color.rgb = ACCENT; q.line.fill.background(); q.shadow.inherit = False
    _cx(s, Inches(0.9), Inches(3.35), Inches(11), Inches(1.4),
        "Demonstrações Financeiras", 46, WHITE, bold=True, font=FONT)
    _cx(s, Inches(0.9), Inches(4.5), Inches(11), Inches(0.7),
        f"{titulo_emp} · {MES_EXT[mes-1]} / {ano}", 24, ACCENT, font=FONT)
    tag = "realizado" if mes <= realizado_ate else "PROJEÇÃO"
    _cx(s, Inches(0.9), Inches(5.2), Inches(11), Inches(0.5),
        f"Fechamento {tag} · realizado até {MESES[realizado_ate-1] if realizado_ate else '—'}/{ano}",
        13, GRAY, font=FONT)
    _cx(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
        "Gerado pelo Cockpit Financeiro Grupo REF · valores em R$", 10, GRAY, font=FONT)
    return s


def _card_kpi(s, x, y, w, h, rotulo, valor, sub, cor_sub=GRAY):
    _retangulo(s, x, y, w, h)
    _cx(s, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.4),
        rotulo, 11, GRAY, bold=True)
    _cx(s, x + Inches(0.25), y + Inches(0.62), w - Inches(0.5), Inches(0.7),
        valor, 27, INK, bold=True)
    _cx(s, x + Inches(0.25), y + h - Inches(0.5), w - Inches(0.5), Inches(0.4),
        sub, 11, cor_sub)


def _slide_kpis(prs, dados):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Destaques do período", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · {MES_EXT[dados['mes']-1]}/{dados['ano']} · variação vs. mesmo período do ano anterior",
        13, GRAY)
    cards = [
        ("RECEITA BRUTA (ANO)", _mi(dados["rb"]), dados["rb_yoy"]),
        ("RECEITA LÍQUIDA (ANO)", _mi(dados["rl"]), dados["rl_yoy"]),
        ("EBIT NEGÓCIO", _pct(dados["ebit_pct"]),
         f"meta {A.META_EBIT_PCT:g}%"),
        ("RESULTADO LÍQUIDO (ANO)", _mi(dados["rliq"]), dados["rliq_yoy"]),
    ]
    x0, y0 = Inches(0.7), Inches(2.0)
    cw, ch, gap = Inches(2.95), Inches(1.9), Inches(0.15)
    for i, (rot, val, sub) in enumerate(cards):
        cor = GRAY
        if isinstance(sub, str) and sub.startswith("▲"): cor = GREEN
        if isinstance(sub, str) and sub.startswith("▼"): cor = RED
        _card_kpi(s, x0 + i * (cw + gap), y0, cw, ch, rot, val, sub, cor)
    # folha/headcount
    _cx(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.5),
        f"Folha do mês ({MESES[dados['mes']-1]}): {_mi(dados['folha'])}   ·   "
        f"Headcount: {dados['headcount']} colaboradores", 14, INK, bold=True)
    return s


def _grafico_barras(s, x, y, w, h, categorias, series, cores, titulo,
                    empilhado=False, rotulos=False, num_fmt='#,##0.0,,'):
    _cx(s, x, y - Inches(0.05), w, Inches(0.4), titulo, 15, INK, bold=True)
    cd = CategoryChartData()
    cd.categories = categorias
    for nome, vals in series:
        cd.add_series(nome, vals)
    tipo = XL_CHART_TYPE.COLUMN_STACKED if empilhado else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = s.shapes.add_chart(tipo, x, y + Inches(0.4), w, h - Inches(0.4), cd)
    ch = gf.chart
    pos = XL_LABEL_POSITION.INSIDE_END if empilhado else XL_LABEL_POSITION.OUTSIDE_END
    _estilo_grafico(ch, cores, com_legenda=len(series) > 1, rotulos=rotulos,
                    rotulo_pos=pos, num_fmt=num_fmt)
    ch.value_axis.has_major_gridlines = True
    ch.value_axis.major_gridlines.format.line.color.rgb = LINE
    ch.value_axis.major_gridlines.format.line.width = Pt(0.5)
    ch.value_axis.format.line.fill.background()
    ch.category_axis.tick_labels.font.size = Pt(9)
    ch.value_axis.tick_labels.number_format = num_fmt
    ch.value_axis.tick_labels.number_format_is_linked = False
    return ch


def _slide_receita_mensal(prs, dados):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Receita bruta mês a mês", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · {dados['ano']} · barras claras a partir de "
        f"{MESES[dados['realizado_ate']-1] if dados['realizado_ate'] else 'Jan'} são projeção",
        13, GRAY)
    ch = _grafico_barras(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.9),
                         MESES, [("Receita Bruta", dados["rb_mensal"])],
                         [ACCENT], "", rotulos=False)
    # esmaece meses projetados (cor por ponto)
    serie = ch.plots[0].series[0]
    for i, pt in enumerate(serie.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = ACCENT if (i + 1) <= dados["realizado_ate"] else LINE
    tot1s = sum(v for i, v in enumerate(dados["rb_mensal"]) if (i + 1) <= dados["realizado_ate"])
    _cx(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
        f"Total realizado até {MESES[dados['realizado_ate']-1] if dados['realizado_ate'] else '—'}: {_mi(tot1s)}",
        13, INK, bold=True)
    return s


def _slide_dre(prs, dados):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Da receita ao resultado (ano)", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · {dados['ano']} · principais linhas da DRE", 13, GRAY)
    cats = ["Rec. Bruta", "Rec. Líquida", "Result. Agência", "EBIT", "Result. Líquido"]
    vals = [dados["rb"], dados["rl"], dados["ra"], dados["ebit"], dados["rliq"]]
    _grafico_barras(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.7),
                    cats, [("R$", vals)], [INK], "", rotulos=True)
    return s


def _slide_trimestral(prs, dados):
    if not dados.get("tri"):
        return None
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Comparativo trimestral", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · Receita bruta por trimestre", 13, GRAY)
    anos = sorted(dados["tri"].keys())
    cats = ["1º Tri", "2º Tri", "3º Tri", "4º Tri"]
    series = [(str(a), [dados["tri"][a].get(t, 0) for t in (1, 2, 3, 4)]) for a in anos]
    cores = [LINE, GRAY, ACCENT][-len(series):] if len(series) <= 3 else [LINE, GRAY, BLUE, ACCENT]
    _grafico_barras(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.9),
                    cats, series, cores, "")
    return s


def _slide_anual(prs, dados):
    if not dados.get("anual"):
        return None
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Evolução anual", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · Receita bruta e resultado líquido por ano", 13, GRAY)
    anos = [str(a["ano"]) for a in dados["anual"]]
    rb = [a["rb"] for a in dados["anual"]]
    _grafico_barras(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.9),
                    anos, [("Receita Bruta", rb)], [ACCENT], "")
    return s


def _slide_despesas(prs, dados):
    if not dados.get("desp_meses"):
        return None
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s)
    _cx(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
        "Composição das despesas", 32, INK, bold=True)
    _cx(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.4),
        f"{dados['titulo']} · {dados['ano']} · pessoal, infraestrutura e outras (R$/mês)", 13, GRAY)
    dm = dados["desp_meses"]
    series = [("Pessoal", [m["pessoal"] for m in dm]),
              ("Infraestrutura", [m["infra"] for m in dm]),
              ("Outras", [m["outras"] for m in dm])]
    _grafico_barras(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.9),
                    MESES, series, [INK, GRAY, ACCENT], "", empilhado=True)
    return s


def _slide_encerramento(prs, dados):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _fundo(s, INK)
    _cx(s, Inches(0.9), Inches(2.9), Inches(11), Inches(1.0),
        "Leitura executiva", 34, WHITE, bold=True)
    linhas = dados.get("insights") or ["—"]
    _cx(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(2.5),
        "\n".join("•  " + t for t in linhas), 16, RGBColor(0xE0, 0xE0, 0xE0))
    _cx(s, Inches(0.9), Inches(6.8), Inches(11), Inches(0.4),
        "Cockpit Financeiro Grupo REF · material interno — não constitui parecer contábil formal",
        10, GRAY)
    return s


# ── coleta de dados (mesma lógica dos endpoints) ─────────────────────────────
def _coleta(slug, ano, mes):
    emp = A.GRUPO if slug == "grupo" else A.BY_SLUG.get(slug)
    if not emp:
        raise ValueError(f"empresa desconhecida: {slug}")
    titulo = "Grupo REF" if slug == "grupo" else emp["label"]
    with A._conn() as con:
        cur = con.cursor(); A._ensure_tables(cur)
        ano = ano or A._ano_default(cur)
        realizado_ate = A._realizado_ate(ano)
        mes = mes or A._folha_periodo_default(cur, ano)

        cur_ = _dre_ano_safe(cur, emp, ano)
        prev = _dre_ano_safe(cur, emp, ano - 1)
        rb, rl = cur_.get(A.RB), cur_.get(A.RL)
        ra, ebit, rliq = cur_.get(A.RA), cur_.get(A.EBIT), cur_.get(A.RLIQ)
        ebit_pct = (ebit / rb * 100) if (ebit is not None and rb) else None

        def yoy(atual, ant):
            if atual is None or not ant:
                return "—"
            d = (atual - ant) / abs(ant) * 100
            seta = "▲" if d >= 0 else "▼"
            return f"{seta} {abs(d):.1f}% vs {ano-1}".replace(".", ",")

        pm = A._dre_mensal(cur, emp, ano)
        rb_mensal = [pm.get(m, {}).get(A.RB, 0) or 0 for m in range(1, 13)]
        folha, headcount = A._folha_mes_total(cur, emp, ano, mes)

        # trimestral (ano corrente por soma dos meses; hist se houver)
        tri = {}
        tri[ano] = {}
        for t in (1, 2, 3, 4):
            meses_t = range((t - 1) * 3 + 1, t * 3 + 1)
            tri[ano][t] = sum(pm.get(m, {}).get(A.RB, 0) or 0 for m in meses_t)
        # anos anteriores via fato_dre_tri_hist (RECEITA_BRUTA por trimestre)
        try:
            ef, ep = A._emp_where(emp)
            cur.execute(f"""SELECT f.ano, f.tri, SUM(f.valor)
                FROM fato_dre_tri_hist f JOIN dim_empresa e ON e.id=f.empresa_id
                WHERE f.metrica='RECEITA_BRUTA'{ef}
                GROUP BY f.ano, f.tri""", ep)
            for a, t, v in cur.fetchall():
                if v is not None:
                    tri.setdefault(int(a), {})[int(t)] = float(v)
        except Exception:
            pass
        # histórico anual
        anual = []
        cur.execute(f"""SELECT p.ano, SUM(CASE WHEN c.descricao=%s THEN f.valor END)
            FROM fato_dre_mensal f JOIN dim_conta c ON c.id=f.conta_id
            JOIN dim_empresa e ON e.id=f.empresa_id JOIN dim_periodo p ON p.id=f.periodo_id
            WHERE 1=1{A._emp_where(emp)[0]}
            GROUP BY p.ano ORDER BY p.ano""",
            [A.RB] + A._emp_where(emp)[1])
        for a, v in cur.fetchall():
            if v:
                anual.append({"ano": int(a), "rb": float(v)})

        # despesas mensais
        desp_meses = [{"pessoal": pm.get(m, {}).get(A.PESSOAL, 0) or 0,
                       "infra": pm.get(m, {}).get(A.INFRA, 0) or 0,
                       "outras": pm.get(m, {}).get(A.OUTRAS, 0) or 0} for m in range(1, 13)]

        insights = []
        if ebit_pct is not None:
            comp = "acima" if ebit_pct >= A.META_EBIT_PCT else "abaixo"
            insights.append(f"EBIT Negócio do ano em {_pct(ebit_pct)} — {comp} da meta de {A.META_EBIT_PCT:g}%.")
        if rb and rliq is not None:
            insights.append(f"Receita bruta acumulada de {_mi(rb)} e resultado líquido de {_mi(rliq)} no ano.")
        melhor = max(range(realizado_ate or 1), key=lambda i: rb_mensal[i]) if realizado_ate else 0
        insights.append(f"Melhor mês realizado: {MES_EXT[melhor]} ({_mi(rb_mensal[melhor])}).")

    return {
        "slug": slug, "titulo": titulo, "ano": ano, "mes": mes,
        "realizado_ate": realizado_ate,
        "rb": rb, "rl": rl, "ra": ra, "ebit": ebit, "rliq": rliq, "ebit_pct": ebit_pct,
        "rb_yoy": yoy(rb, prev.get(A.RB)), "rl_yoy": yoy(rl, prev.get(A.RL)),
        "rliq_yoy": yoy(rliq, prev.get(A.RLIQ)),
        "rb_mensal": rb_mensal, "folha": folha, "headcount": headcount,
        "tri": tri if len(tri) else None, "anual": anual, "desp_meses": desp_meses,
        "insights": insights,
    }


def _dre_ano_safe(cur, emp, ano):
    try:
        return A._dre_ano(cur, emp, ano)
    except Exception:
        return {}


# ── API pública ──────────────────────────────────────────────────────────────
def gerar_pptx(slug="grupo", ano=None, mes=None):
    """Gera o deck e devolve (bytes_pptx, nome_arquivo)."""
    d = _coleta(slug, ano, mes)
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    _slide_capa(prs, d["titulo"], d["ano"], d["mes"], d["realizado_ate"])
    _slide_kpis(prs, d)
    _slide_receita_mensal(prs, d)
    _slide_dre(prs, d)
    _slide_trimestral(prs, d)
    _slide_anual(prs, d)
    _slide_despesas(prs, d)
    _slide_encerramento(prs, d)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    nome = f"Demonstracoes_Financeiras_{d['titulo'].replace(' ', '_')}_{d['ano']}_{MESES[d['mes']-1]}.pptx"
    return buf.read(), nome


def run(slug="grupo", ano=None, mes=None):
    """Para o runner: grava em data/relatorios/ e devolve {ok, output, path}."""
    try:
        blob, nome = gerar_pptx(slug, ano, mes)
        outdir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "data", "relatorios")
        os.makedirs(outdir, exist_ok=True)
        path = os.path.join(outdir, nome)
        with open(path, "wb") as f:
            f.write(blob)
        return {"ok": True, "output": f"relatório gerado: {nome} ({len(blob)} bytes)", "path": path}
    except Exception as e:
        return {"ok": False, "output": f"ERRO relatório: {e}"}


def main():
    slug = sys.argv[1] if len(sys.argv) > 1 else "grupo"
    ano = int(sys.argv[2]) if len(sys.argv) > 2 else None
    mes = int(sys.argv[3]) if len(sys.argv) > 3 else None
    r = run(slug, ano, mes)
    print(r["output"])
    sys.exit(0 if r["ok"] else 1)


if __name__ == "__main__":
    main()
